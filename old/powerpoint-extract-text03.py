import streamlit as st
import pandas as pd
from pptx import Presentation
import requests
import json
from io import BytesIO
import zipfile

# VOICEVOXのAPI設定
host = "127.0.0.1"
port = 50021

def generate_audio(text, speaker):
    params = {"text": text, "speaker": speaker}
    response1 = requests.post(f"http://{host}:{port}/audio_query", params=params)
    response2 = requests.post(
        f"http://{host}:{port}/synthesis",
        headers={"Content-Type": "application/json"},
        params=params,
        data=json.dumps(response1.json())
    )
    return response2.content

def extract_text_from_pptx(uploaded_file, speaker):
    prs = Presentation(uploaded_file)
    text_data = []
    audio_files = {}

    for i, slide in enumerate(prs.slides):
        notes_slide = slide.notes_slide
        text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
        text_data.append({"slide": i + 1, "text": text})
        if text:
            audio_files[f"slide_{i:02}.wav"] = generate_audio(text, speaker)

    return pd.DataFrame(text_data), audio_files

st.title("PowerPoint Notes to Audio Converter")

uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
if uploaded_file:
    speaker = st.selectbox("Choose a speaker voice", [1, 2, 3, 4])  # Speaker IDの選択肢
    start_button = st.button("Start TTS")

    if start_button:
        text_data, audio_files = extract_text_from_pptx(uploaded_file, speaker)
        st.write(text_data)

        if audio_files:
            # ファイルをZIP形式でダウンロード
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for filename, content in audio_files.items():
                    zip_file.writestr(filename, content)
            st.download_button(label="Download Audio Files",
                               data=zip_buffer.getvalue(),
                               file_name="audio_files.zip",
                               mime="application/zip")
            # 使用後に辞書をクリアしてメモリ解放
            audio_files.clear()
