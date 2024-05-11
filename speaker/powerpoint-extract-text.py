import requests
import json
from pptx import Presentation

# VOICEVOXの設定
host = "127.0.0.1"
port = 50021
params = {
    "speaker": 2,  # 音声の種類をInt型で指定
}

# PowerPointファイルの読み込み
ppt_file_path = "../samplepowerpoint.pptx"
prs = Presentation(ppt_file_path)

# 各スライドのノートを読み上げて音声ファイルとして保存
for i, slide in enumerate(prs.slides):
    notes_slide = slide.notes_slide
    text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""

    if text:  # ノートにテキストがある場合のみ処理
        response1 = requests.post(
            f"http://{host}:{port}/audio_query",
            params={**params, "text": text}
        )
        response2 = requests.post(
            f"http://{host}:{port}/synthesis",
            headers={"Content-Type": "application/json"},
            params=params,
            data=json.dumps(response1.json())
        )

        # 音声ファイルをスライド番号に基づいて保存
        audio_file_path = f"slide_{i:02}.wav"
        with open(audio_file_path, "wb") as f:
            f.write(response2.content)
